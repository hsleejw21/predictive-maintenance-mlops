# -*- coding: utf-8 -*-
"""기말 발표 PPT 생성 — 설비 고장 모드 분류 + SPC 이중 레이어 MLOps (13장)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pathlib import Path

# 이미지(shots/)와 출력 pptx를 스크립트 위치 기준으로 해석 → 어디서 실행하든 동일
HERE = Path(__file__).resolve().parent

# ---------------- 색상 ----------------
NAVY=RGBColor(0x1e,0x29,0x3b); BLUE=RGBColor(0x25,0x63,0xeb); BLUE_L=RGBColor(0xdb,0xea,0xfe)
AMBER=RGBColor(0xd9,0x77,0x06); AMBER_L=RGBColor(0xfe,0xf3,0xc7)
PURPLE=RGBColor(0x7c,0x3a,0xed); PURPLE_L=RGBColor(0xed,0xe9,0xfe)
GREEN=RGBColor(0x05,0x96,0x69); GREEN_L=RGBColor(0xd1,0xfa,0xe5)
GREEN2=RGBColor(0x34,0xd3,0x99); SLATE=RGBColor(0xcb,0xd5,0xe1)
RED=RGBColor(0xdc,0x26,0x26); RED_L=RGBColor(0xfe,0xe2,0xe2)
GRAY=RGBColor(0x6b,0x72,0x80); GRAY_L=RGBColor(0xf3,0xf4,0xf6)
WHITE=RGBColor(0xff,0xff,0xff); DARK=RGBColor(0x1f,0x29,0x37); LINE=RGBColor(0xe5,0xe7,0xeb)
FONT="Apple SD Gothic Neo"

prs=Presentation()
prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide,x,y,w,h,fill,line=None,lw=1.0,rounded=False,shadow=False):
    shp=slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x),Inches(y),Inches(w),Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    if rounded:
        try: shp.adjustments[0]=0.06
        except Exception: pass
    return shp

def txt(slide,x,y,w,h,lines,size=14,color=DARK,bold=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,font=FONT,space_after=4,line_spacing=1.05):
    """lines: str 또는 [(text,size,color,bold,align,bullet_level)] 리스트(부분 dict 허용)."""
    tb=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    tf.margin_top=Inches(0.02); tf.margin_bottom=Inches(0.02)
    if isinstance(lines,str): lines=[lines]
    first=True
    for ln in lines:
        if isinstance(ln,dict):
            t=ln.get("t",""); s=ln.get("size",size); c=ln.get("color",color)
            b=ln.get("bold",bold); a=ln.get("align",align); lvl=ln.get("level",0)
            sa=ln.get("space_after",space_after); ls=ln.get("line_spacing",line_spacing)
        else:
            t=ln; s=size; c=color; b=bold; a=align; lvl=0; sa=space_after; ls=line_spacing
        p=tf.paragraphs[0] if first else tf.add_paragraph()
        first=False
        p.alignment=a; p.space_after=Pt(sa); p.space_before=Pt(0)
        try: p.line_spacing=ls
        except Exception: pass
        p.level=lvl
        r=p.add_run(); r.text=t
        r.font.size=Pt(s); r.font.bold=b; r.font.color.rgb=c; r.font.name=font
    return tb

def title_bar(slide,title,kicker=None,accent=BLUE):
    rect(slide,0,0,SW.inches,1.12,NAVY)
    rect(slide,0,1.12,SW.inches,0.06,accent)
    if kicker:
        txt(slide,0.6,0.12,11,0.3,kicker,size=12,color=BLUE_L,bold=True)
        txt(slide,0.6,0.40,12.1,0.7,title,size=25,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    else:
        txt(slide,0.6,0.1,12.1,1.0,title,size=26,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)

TOTAL=16
_pgctr=[1]   # 표지=1 (번호 미표시). pagenum 호출마다 자동 증가
def pagenum(slide,n=None):
    _pgctr[0]+=1
    txt(slide,12.3,7.05,0.9,0.35,f"{_pgctr[0]} / {TOTAL}",size=10,color=GRAY,align=PP_ALIGN.RIGHT)

def add_img(slide,path,x,y,w):
    p=Path(path)
    if not p.is_absolute(): p=HERE/p          # 상대경로는 스크립트 기준
    pic=slide.shapes.add_picture(str(p),Inches(x),Inches(y),width=Inches(w))
    pic.line.color.rgb=LINE; pic.line.width=Pt(1.2)
    return pic

def table(slide,x,y,w,data,col_w,header=True,fontsize=12.5,row_h=0.42,
          head_fill=NAVY,head_color=WHITE,cell_colors=None):
    """Canva 친화형: 네이티브 표 대신 셀마다 도형+텍스트로 그려 편집 가능하게.
    data: 행 리스트(셀=str 또는 (text,color)). cell_colors: {(r,c):RGB}."""
    rows=len(data); cols=len(data[0]); tot=sum(col_w)
    xs=[x]
    for j in range(cols): xs.append(xs[-1]+w*col_w[j]/tot)
    for i in range(rows):
        cy=y+i*row_h
        for j in range(cols):
            cx=xs[j]; cwj=xs[j+1]-xs[j]
            val=data[i][j]; color=DARK
            if isinstance(val,tuple): val,color=val[0],val[1]
            if header and i==0:
                fill=head_fill; tcolor=head_color; bold=True
            else:
                fill=WHITE; tcolor=color; bold=(color in (RED,GREEN))
                if cell_colors and (i,j) in cell_colors: fill=cell_colors[(i,j)]
            rect(slide,cx,cy,cwj,row_h,fill,line=LINE,lw=0.75)
            txt(slide,cx+0.12,cy,cwj-0.22,row_h,str(val),size=fontsize,color=tcolor,
                bold=bold,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.0)
    return None

def arrow_down(slide,x,y,label=None):
    txt(slide,x,y,1.2,0.3,"▼",size=16,color=GRAY,align=PP_ALIGN.CENTER)
    if label:
        txt(slide,x+0.9,y+0.0,3.0,0.3,label,size=11,color=GRAY)

# ============================================================ Slide 1 표지
s=blank()
rect(s,0,0,SW.inches,SH.inches,NAVY)
rect(s,0,5.0,SW.inches,0.06,BLUE)
rect(s,0.0,0.0,0.35,SH.inches,BLUE)
txt(s,1.0,1.7,11.3,0.4,"데이터애널리틱스특론 · 기말 프로젝트",size=15,color=BLUE_L,bold=True)
txt(s,1.0,2.2,11.3,2.0,[
    {"t":"설비 고장 모드 분류 + SPC 관리도","size":40,"color":WHITE,"bold":True,"space_after":6},
    {"t":"이중 레이어 MLOps 시스템","size":40,"color":WHITE,"bold":True},
])
txt(s,1.0,4.15,11.3,0.6,"UCI AI4I 2020 Predictive Maintenance 데이터셋 기반 시연",
    size=17,color=RGBColor(0x94,0xa3,0xb8))
txt(s,1.0,5.4,11.3,1.2,[
    {"t":"문제정의  →  MLOps로 해결  →  기대효과","size":16,"color":BLUE_L,"bold":True,"space_after":10},
    {"t":"발표자: ___________      2026","size":14,"color":RGBColor(0x94,0xa3,0xb8)},
])

# ============================================================ Slide 2 목차
s=blank(); title_bar(s,"목차","CONTENTS")
agenda=[
    ("1  배경", BLUE, ["문제 정의 — 현장의 4가지 Pain (오탐·원인불명·드리프트·라벨지연)",
                       "왜 MLOps인가 — 한 번 만든 모델이 늙는다"]),
    ("2  솔루션", PURPLE, ["이중 레이어 아키텍처 개요",
                          "Layer 1 · SPC 관리도 (Shewhart 3σ)",
                          "Layer 2 · Random Forest 다중분류",
                          "이중 레이어의 핵심 / MLOps 전체 루프"]),
    ("3  검증", AMBER, ["데이터셋 (UCI AI4I 2020)",
                        "드리프트 시나리오 (베어링 마모)",
                        "기대 효과 — 구현 실측 결과"]),
    ("4  데모 & 결론", GREEN, ["실시간 대시보드 (Streamlit) 시연",
                              "결론 & 향후 계획"]),
]
y=1.5
for atitle,c,items in agenda:
    h=0.32*len(items)+0.42
    rect(s,0.6,y,2.9,h,c,rounded=True)
    txt(s,0.6,y,2.9,h,atitle,size=18,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,3.7,y,9.03,h,GRAY_L,rounded=True)
    txt(s,4.0,y,8.5,h,[{"t":"·  "+it,"size":13.5,"space_after":4,"line_spacing":1.1} for it in items],
        anchor=MSO_ANCHOR.MIDDLE)
    y+=h+0.12
pagenum(s)

# ============================================================ Slide 3 한 장 요약
s=blank(); title_bar(s,"한 장 요약","EXECUTIVE SUMMARY")
cards=[("무엇을",BLUE,BLUE_L,
        "제조 설비의 5개 센서로 고장을 실시간 감지하고, 5개 고장 모드를 자동 진단하는 MLOps 시스템"),
       ("왜",RED,RED_L,
        "단순 임계값은 오탐이 많고 원인 진단이 안 됨. 한 번 만든 모델은 설비 변화(드리프트)로 늙음"),
       ("어떻게",GREEN,GREEN_L,
        "SPC 관리도(1차 감시) + Random Forest(2차 진단) 이중 레이어 + 자동 재학습 루프")]
cw=3.9
for i,(t,c,cl,body) in enumerate(cards):
    x=0.6+i*4.05
    rect(s,x,1.6,cw,3.4,cl,rounded=True)
    rect(s,x,1.6,cw,0.7,c,rounded=True)
    txt(s,x,1.6,cw,0.7,t,size=20,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.25,2.55,cw-0.5,2.3,body,size=15,color=DARK,line_spacing=1.25)
rect(s,0.6,5.35,12.13,1.2,NAVY,rounded=True)
txt(s,0.9,5.35,11.5,1.2,[
    {"t":"핵심 결과 (예고)","size":13,"color":BLUE_L,"bold":True,"space_after":4},
    {"t":"드리프트 구간에서 — 재학습을 안 하면 고장 탐지 성능이 급락하지만, 이중 레이어가 드리프트를 자동 감지·재학습해 성능을 회복한다.",
     "size":15,"color":WHITE,"bold":True}],anchor=MSO_ANCHOR.MIDDLE)
pagenum(s,2)

# ============================================================ Slide 3 문제 정의
s=blank(); title_bar(s,"문제 정의 — 현장은 무엇 때문에 힘든가","PROBLEM",accent=RED)
# Scene
rect(s,0.6,1.55,5.7,2.35,GRAY_L,rounded=True)
txt(s,0.85,1.7,5.3,0.4,"Scene — 어디서",size=15,color=NAVY,bold=True)
txt(s,0.85,2.15,5.3,1.7,[
    {"t":"회전·절삭 설비(밀링머신 등)가 24시간 가동.","size":13.5,"space_after":6},
    {"t":"5개 센서가 가동 내내 측정:","size":13.5,"bold":True,"space_after":4},
    {"t":"공기온도 · 공정온도 · 회전속도 · 토크 · 공구마모","size":13.5,"color":BLUE,"bold":True},
])
# Pain
rect(s,6.5,1.55,6.23,2.35,RED_L,rounded=True)
txt(s,6.75,1.7,5.8,0.4,"Pain — 무엇이 힘든가",size=15,color=RED,bold=True)
txt(s,6.75,2.15,5.8,1.7,[
    {"t":"• 단순 임계값 알람 → 오탐 多, 복합 고장은 놓침","size":12.5,"space_after":5},
    {"t":"• 「고장」은 알아도 「어떤 고장」인지 모름 → 매번 분해","size":12.5,"space_after":5},
    {"t":"• 설비 노후·부품교체로 정상 패턴 변화 → 모델이 늙음","size":12.5,"space_after":5},
    {"t":"• 실제 원인 라벨은 정비 후 확정 → 라벨 지연","size":12.5},
])
# KPI
rect(s,0.6,4.1,12.13,2.4,NAVY,rounded=True)
txt(s,0.9,4.25,11.5,0.4,"KPI — 무엇을 개선하고 싶은가",size=15,color=BLUE_L,bold=True)
kpis=["비계획 다운타임 ↓","고장 원인 진단 시간 ↓","오탐 알람 수 ↓","모델 유지보수 공수 ↓"]
for i,k in enumerate(kpis):
    x=0.95+i*2.95
    rect(s,x,4.85,2.7,1.35,RGBColor(0x33,0x41,0x55),rounded=True)
    txt(s,x,4.85,2.7,1.35,k,size=15,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
pagenum(s,3)

# ============================================================ Slide 4 왜 MLOps
s=blank(); title_bar(s,"왜 MLOps인가 — 일반 ML의 한계","WHY MLOPS")
txt(s,0.6,1.45,12.1,0.4,"\"모델을 만들었다 = 끝났다\"가 아니다. 설비 환경은 계속 변하므로 모델은 배포 순간부터 늙기 시작한다.",
    size=15,color=DARK,bold=True)
rect(s,0.6,2.0,12.13,1.5,GRAY_L,rounded=True)
txt(s,0.85,2.1,11.6,0.35,"일반 ML 프로젝트의 전형적 비극",size=13,color=NAVY,bold=True)
txt(s,0.85,2.5,11.8,1.0,
    "① 6개월 들여 정확도 95% 달성  →  ② 수동 배포  →  ③ 2개월 뒤 80%로 하락  "
    "→  ④ 원인·데이터·코드 행방 불명  →  ⑤ 처음부터 다시 6개월",
    size=13.5,color=DARK,line_spacing=1.2)
table(s,0.6,3.8,12.13,[
    ["항목","일반 ML (1회 PoC)","MLOps"],
    ["배포",("엑셀·리포트로 결과 전달",RED),("실시간 추론으로 자동 서빙",GREEN)],
    ["성능 하락 감지",("현장 컴플레인이 와야 인지",RED),("SPC + PSI로 24시간 자동 감시",GREEN)],
    ["재학습",("사람이 처음부터 다시 (수주)",RED),("드리프트 트리거 → 자동 재학습·배포",GREEN)],
],[2.5,4.8,4.8],row_h=0.62,fontsize=13.5)
pagenum(s,4)

# ============================================================ Slide 5 아키텍처
s=blank(); title_bar(s,"해결 — 이중 레이어 아키텍처","SOLUTION OVERVIEW",accent=GREEN)
cx=3.4; cw=6.5
def archbox(y,h,fill,line,title,sub,tag):
    rect(s,cx,y,cw,h,fill,line=line,lw=1.5,rounded=True)
    txt(s,cx+0.25,y+0.1,cw-0.5,0.4,title,size=15,color=line,bold=True)
    txt(s,cx+0.25,y+0.5,cw-0.5,h-0.55,sub,size=12.5,color=DARK,line_spacing=1.12)
    if tag:
        txt(s,cx+cw+0.15,y+0.1,3.4,h,tag,size=11.5,color=GRAY,anchor=MSO_ANCHOR.MIDDLE)
txt(s,cx,1.45,cw,0.4,"📡 센서 스트림  (AI4I 2020, 행 단위 주입)",size=14,color=NAVY,bold=True,align=PP_ALIGN.CENTER)
arrow_down(s,cx+cw/2-0.6,1.85)
archbox(2.2,1.05,AMBER_L,AMBER,"Layer 1 — SPC 관리도 (Shewhart 3σ)",
        "5개 센서 각각 관리한계(μ±3σ) 이탈 시 즉시 알람","실시간 1차 감시\n라벨 불필요·해석 명확")
arrow_down(s,cx+cw/2-0.6,3.3,"이상 플래그")
archbox(3.65,1.05,PURPLE_L,PURPLE,"Layer 2 — Random Forest 다중분류",
        "플래그된 시점을 5개 고장 모드로 분류 + 「왜」 설명","2차 정밀 진단")
arrow_down(s,cx+cw/2-0.6,4.75)
archbox(5.1,1.15,RED_L,RED,"교차 검증 → 자동 재학습 트리거",
        "SPC=이상인데 RF=정상/낮은확신 → 신종 패턴(컨셉 드리프트) → 재학습","두 레이어가\n서로를 감시")
pagenum(s,5)

# ============================================================ Slide 6 Layer1 SPC
s=blank(); title_bar(s,"Layer 1 — SPC 관리도 (Shewhart 3σ)","LAYER 1 · MONITORING",accent=AMBER)
rect(s,0.6,1.5,5.9,2.3,GRAY_L,rounded=True)
txt(s,0.85,1.6,5.4,0.4,"작동 방식 — 개별값 관리도(I-chart)",size=14,color=NAVY,bold=True)
txt(s,0.85,2.05,5.4,1.7,[
    {"t":"• 기준기간 데이터로 센서별 중심선 μ 계산","size":13,"space_after":5},
    {"t":"• 관리상한 UCL = μ + 3σ","size":13,"color":AMBER,"bold":True,"space_after":3},
    {"t":"• 관리하한 LCL = μ − 3σ","size":13,"color":AMBER,"bold":True,"space_after":5},
    {"t":"• 한계 이탈(OOC) → 즉시 알람 (오탐 ≈0.27%)","size":13},
])
# 간이 관리도 그림
gx,gy,gw,gh=6.8,1.5,5.9,2.3
rect(s,gx,gy,gw,gh,WHITE,line=LINE,lw=1)
# UCL/center/LCL 선 (Canva 친화: 커넥터 대신 얇은 사각형)
for yy,lab,col in [(gy+0.5,"UCL",RED),(gy+gh/2,"center",GRAY),(gy+gh-0.6,"LCL",RED)]:
    rect(s,gx+0.2,yy-0.012,gw-1.1,0.024,col)
    txt(s,gx+gw-0.75,yy-0.12,0.7,0.25,lab,size=9,color=col)
# 점들 (정상/이탈)
import random; random.seed(1)
pts=[gy+gh/2+random.uniform(-0.35,0.35) for _ in range(11)]+[gy+0.35,gy+0.3]
for i,py in enumerate(pts):
    px=gx+0.35+i*0.42
    ooc = py<gy+0.5
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,Inches(px),Inches(py),Inches(0.12),Inches(0.12))
    dot.fill.solid(); dot.fill.fore_color.rgb=RED if ooc else BLUE; dot.line.fill.background()
txt(s,gx+0.1,gy+gh-0.35,gw,0.3,"빨강 = 관리이탈(OOC) → Layer-1 경보",size=10,color=GRAY)
# 두 역할
rect(s,0.6,4.05,12.13,2.5,AMBER_L,rounded=True)
txt(s,0.85,4.18,11.6,0.4,"SPC가 맡는 두 가지 역할 (이 설계의 묘미)",size=15,color=AMBER,bold=True)
rect(s,0.95,4.7,5.7,1.6,WHITE,rounded=True)
txt(s,1.2,4.82,5.2,1.4,[
    {"t":"① 실시간 이상 감지","size":14,"color":NAVY,"bold":True,"space_after":4},
    {"t":"라벨 없이 즉시 「평소와 다르다」 포착. RF가 학습한 적 없는 상황에서도 작동(cold-start 강건).","size":12.5,"line_spacing":1.15},
])
rect(s,6.75,4.7,5.7,1.6,WHITE,rounded=True)
txt(s,7.0,4.82,5.2,1.4,[
    {"t":"② ML 모델의 입력 드리프트 감시","size":14,"color":NAVY,"bold":True,"space_after":4},
    {"t":"센서 분포가 관리한계를 점점 벗어남 = 입력 드리프트 신호. SPC가 곧 MLOps 모니터링 레이어.","size":12.5,"line_spacing":1.15},
])
pagenum(s,6)

# ============================================================ Slide 7 Layer2 RF
s=blank(); title_bar(s,"Layer 2 — Random Forest 다중분류","LAYER 2 · DIAGNOSIS",accent=PURPLE)
table(s,0.6,1.5,12.13,[
    ["코드","고장 모드","물리적 원인 (AI4I 공식 규칙)"],
    ["TWF","공구마모 고장","공구 마모 200~240분 구간에서 무작위 발생"],
    ["HDF","열발산 고장","(공정온도−공기온도) < 8.6K 이면서 회전속도 < 1380rpm"],
    ["PWF","전력 고장","일률(토크×각속도) < 3500W 또는 > 9000W"],
    ["OSF","과부하 고장","공구마모×토크 > 임계값 (L/M/H = 11000/12000/13000)"],
    ["RNF","무작위 고장","0.1% 확률로 무작위 발생"],
],[1.3,2.6,8.2],row_h=0.46,fontsize=12.5,head_fill=PURPLE)
rect(s,0.6,4.95,5.95,1.5,GRAY_L,rounded=True)
txt(s,0.85,5.05,5.5,1.3,[
    {"t":"입력 특징","size":13.5,"color":NAVY,"bold":True,"space_after":4},
    {"t":"제품 타입(L/M/H) + 5개 센서 + 파생특징(temp_diff, power) — 물리 규칙을 모델이 쉽게 학습하도록","size":12,"line_spacing":1.15},
])
rect(s,6.78,4.95,5.95,1.5,GRAY_L,rounded=True)
txt(s,7.03,5.05,5.5,1.3,[
    {"t":"설계 포인트","size":13.5,"color":NAVY,"bold":True,"space_after":4},
    {"t":"• 클래스 불균형(고장률 3.4%) → class_weight=balanced","size":12,"space_after":3},
    {"t":"• feature importance로 「왜 이 고장」 근거 제시","size":12},
])
pagenum(s,7)

# ============================================================ Slide 8 이중레이어 핵심
s=blank(); title_bar(s,"이중 레이어의 핵심 — 두 레이어가 서로 검증","KEY IDEA",accent=PURPLE)
txt(s,0.6,1.4,12.1,0.4,"SPC × RF 결과를 교차하면 4가지 상황. 특히 「SPC 이상 + RF 정상」이 컨셉 드리프트를 잡는 황금 신호.",
    size=14,color=DARK,bold=True)
cells=[("SPC 정상 · RF 정상",GREEN_L,GREEN,"정상 가동. 조치 없음."),
       ("SPC 이상 · RF 고장(높은 확신)",BLUE_L,BLUE,"진단 확정. 어떤 고장인지까지 즉시 정비팀 통보."),
       ("SPC 정상 · RF 고장",AMBER_L,AMBER,"미세 전조. 관리한계엔 안 닿았지만 RF가 포착 → 예의주시."),
       ("★ SPC 이상 · RF 정상/낮은확신",RED_L,RED,"신종 패턴 의심 → 컨셉 드리프트 → 사람 검토 + 재학습 트리거.")]
for i,(t,cl,c,body) in enumerate(cells):
    x=0.6+(i%2)*6.15; y=2.05+(i//2)*1.95
    rect(s,x,y,5.9,1.8,cl,line=c,lw=1.5,rounded=True)
    txt(s,x+0.25,y+0.15,5.4,0.5,t,size=14.5,color=c,bold=True)
    txt(s,x+0.25,y+0.7,5.4,1.0,body,size=13,color=DARK,line_spacing=1.15)
rect(s,0.6,6.15,12.13,0.95,NAVY,rounded=True)
txt(s,0.9,6.15,11.5,0.95,
    "빠르지만 단순한 SPC가 「이상」을 먼저 외치고, 정밀하지만 과거에 갇힌 RF가 「모르겠다」면 → 그 간극이 모델이 따라잡을 새 현실. SPC가 RF의 감시자가 된다.",
    size=13.5,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.15)
pagenum(s,8)

# ============================================================ Slide 9 MLOps 루프
s=blank(); title_bar(s,"MLOps 전체 루프 (5개 레이어)","MLOPS LOOP",accent=GREEN)
layers=[("Data",BLUE,"센서 스트림 → Feature Store (temp_diff·power 파생값을 학습·운영이 동일하게 사용)"),
        ("Model",PURPLE,"Random Forest 학습 + 버전 관리(v1, v2…) + 재현성"),
        ("Serving",AMBER,"실시간 추론 (행 단위 스트림 처리)"),
        ("Monitoring",RED,"SPC 관리도(즉시 이상) + PSI(분포 드리프트) — 이중 감시"),
        ("Automation",GREEN,"PSI 임계 초과 → 자동 재학습 → 검증 → 배포 → 실패 시 롤백")]
for i,(t,c,body) in enumerate(layers):
    y=1.55+i*0.72
    rect(s,0.6,y,2.2,0.62,c,rounded=True)
    txt(s,0.6,y,2.2,0.62,t,size=14,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    rect(s,2.95,y,4.4,0.62,GRAY_L,rounded=True)
    txt(s,3.15,y,4.05,0.62,body,size=11.5,color=DARK,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05)
# 자동 재학습 흐름
rect(s,7.6,1.55,5.13,4.32,PURPLE_L,rounded=True)
txt(s,7.85,1.7,4.6,0.4,"자동 재학습 흐름",size=14,color=PURPLE,bold=True)
steps=["① 드리프트 진입 → 토크 분포 변화","② 롤링 PSI가 0.2 임계 초과",
       "③ 재학습 트리거 → 최근 윈도우로 RF v2 학습","④ v1과 비교 검증 → 배포 또는 롤백",
       "⑤ 「재학습 안 했다면」 반례와 비교 → 가치 입증"]
txt(s,7.85,2.2,4.6,3.5,[{"t":st,"size":12.5,"space_after":11,"line_spacing":1.1} for st in steps])
pagenum(s,9)

# ============================================================ Slide 10 데이터셋
s=blank(); title_bar(s,"데이터셋 — UCI AI4I 2020 Predictive Maintenance","DATASET")
rect(s,0.6,1.55,6.0,3.0,GRAY_L,rounded=True)
txt(s,0.85,1.7,5.5,2.8,[
    {"t":"개요","size":14,"color":NAVY,"bold":True,"space_after":6},
    {"t":"• 10,000건 · 5개 센서 · 5개 고장 모드","size":13,"space_after":5},
    {"t":"• 제품 타입 L / M / H (품질 등급)","size":13,"space_after":5},
    {"t":"• 전체 고장률 3.4% (339건) — 현실적 불균형","size":13,"space_after":5},
    {"t":"• 센서: 공기온도·공정온도·회전속도·토크·공구마모","size":13},
])
# 고장모드별 건수 막대
rect(s,6.8,1.55,5.93,3.0,WHITE,line=LINE,lw=1,rounded=True)
txt(s,7.05,1.65,5.4,0.35,"고장 모드별 건수",size=13.5,color=NAVY,bold=True)
counts=[("HDF",115),("OSF",98),("PWF",95),("TWF",46),("RNF",19)]
maxc=115
for i,(name,cnt) in enumerate(counts):
    y=2.15+i*0.46
    txt(s,7.05,y,0.9,0.35,name,size=12,color=DARK,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    bw=3.6*cnt/maxc
    rect(s,8.0,y+0.04,bw,0.28,PURPLE)
    txt(s,8.05+bw,y,0.8,0.35,str(cnt),size=11,color=GRAY,anchor=MSO_ANCHOR.MIDDLE)
rect(s,0.6,4.75,12.13,1.6,GREEN_L,line=GREEN,lw=1.5,rounded=True)
txt(s,0.9,4.9,11.5,1.4,[
    {"t":"✓ 물리 규칙으로 라벨 100% 재현 검증됨","size":15,"color":GREEN,"bold":True,"space_after":5},
    {"t":"HDF·PWF·OSF를 공식 문서의 결정론적 규칙으로 재계산한 결과 원본 라벨과 완전히 일치(115/115, 95/95, 98/98). → 드리프트를 주입하고 라벨을 재계산해도 물리적으로 정직한 시나리오가 성립한다.",
     "size":13,"color":DARK,"line_spacing":1.2},
])
pagenum(s,10)

# ============================================================ Slide 11 드리프트 시나리오
s=blank(); title_bar(s,"드리프트 시나리오 — \"베어링 마모로 부하가 점점 증가\"","DRIFT SCENARIO",accent=AMBER)
# 타임라인
phases=[("기준기간 baseline","0 ~ 3,000","SPC 한계 산출\n+ RF v1 학습\n+ PSI 기준 확정",BLUE),
        ("안정 가동 stable","3,000 ~ 6,000","정상 분포\n모델 잘 작동",GREEN),
        ("드리프트 drift","6,000 ~ 10,000","토크 +16Nm 점진 주입\n→ 물리 규칙으로\nOSF·PWF 재계산",RED)]
for i,(t,rng,body,c) in enumerate(phases):
    x=0.6+i*4.05
    rect(s,x,1.55,3.85,1.9,WHITE,line=c,lw=2,rounded=True)
    rect(s,x,1.55,3.85,0.5,c,rounded=True)
    txt(s,x,1.55,3.85,0.5,t,size=13,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.2,2.1,3.5,0.35,rng,size=12,color=c,bold=True,align=PP_ALIGN.CENTER)
    txt(s,x+0.2,2.5,3.5,0.9,body,size=12,color=DARK,align=PP_ALIGN.CENTER,line_spacing=1.1)
    if i<2: txt(s,x+3.8,2.3,0.4,0.4,"→",size=20,color=GRAY,align=PP_ALIGN.CENTER)
rect(s,0.6,3.75,12.13,2.75,AMBER_L,rounded=True)
txt(s,0.85,3.9,11.6,0.4,"🔥 연쇄 반응 — 이게 시연의 하이라이트",size=15,color=AMBER,bold=True)
chain=[("①","SPC 토크 관리도가 UCL 이탈 (빨간 점 급증)"),
       ("②","PSI(토크)가 0.2 임계 초과"),
       ("③","과부하(OSF)·전력(PWF) 고장 실제 증가"),
       ("④","기준기간 학습한 RF v1의 고장 탐지 성능 하락"),
       ("⑤","재학습 트리거 → RF v2 학습"),
       ("⑥","성능 회복 (재학습의 가치 입증)")]
for i,(n,body) in enumerate(chain):
    x=0.95+(i%3)*3.95; y=4.45+(i//3)*0.95
    rect(s,x,y,3.7,0.78,WHITE,rounded=True)
    txt(s,x+0.15,y,0.5,0.78,n,size=17,color=AMBER,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+0.65,y,2.9,0.78,body,size=11.5,color=DARK,anchor=MSO_ANCHOR.MIDDLE,line_spacing=1.05)
pagenum(s,11)

# ============================================================ Slide 12 기대효과
s=blank(); title_bar(s,"기대 효과 — 구현 실측 결과","EXPECTED VALUE",accent=GREEN)
txt(s,0.6,1.32,12.1,0.35,"정량 — 드리프트 구간(6,000~10,000행) 재학습 전/후 비교 · AI4I 2020 시뮬레이션 실측",size=13.5,color=NAVY,bold=True)
table(s,0.6,1.75,12.13,[
    ["지표","재학습 안 함 (v1 고정)","이중 레이어 + 자동 재학습"],
    ["고장 탐지 recall (놓치지 않은 고장 비율)",("78.8 %",RED),("90.4 %   (+11.6%p)",GREEN)],
    ["정확도 (accuracy)",("92.1 %",RED),("95.9 %   (+3.8%p)",GREEN)],
    ["종합 성능 (macro-F1)",("0.561",RED),("0.601",GREEN)],
],[5.0,3.4,3.7],row_h=0.58,fontsize=13.5)
rect(s,0.6,4.4,12.13,0.82,NAVY,rounded=True)
txt(s,0.9,4.4,11.5,0.82,"★ 재학습으로 드리프트 구간 고장 탐지율 78.8% → 90.4% 회복 — 재학습이 없었다면 급증한 고장을 다수 놓쳤을 것",
    size=14.5,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
rect(s,0.6,5.45,12.13,1.1,GREEN_L,rounded=True)
txt(s,0.85,5.55,11.6,0.4,"정성 — 현장이 얻는 것",size=13.5,color=GREEN,bold=True)
txt(s,0.85,5.95,11.6,0.55,[
    {"t":"• 감지 + 진단 + 자기유지를 한 시스템이 동시에 수행      • 단순 임계값 대비 복합 고장 조기 포착, 오탐↓","size":12.5,"space_after":4},
    {"t":"• 드리프트 자동 대응 → 모델 유지보수 공수↓      • SPC라는 현장 친화적·해석 가능한 1차 방어선 → 도입 저항 낮음","size":12.5},
])
pagenum(s,12)

# ============================================================ Slide 13 데모 핵심 (money shot)
s=blank(); title_bar(s,"데모 — 드리프트 감지 → 자동 재학습 → 회복","LIVE DEMO",accent=PURPLE)
add_img(s,"shots/04_drift_retrain.png",0.5,1.45,8.75)
rect(s,9.5,1.45,3.25,5.05,NAVY,rounded=True)
txt(s,9.75,1.65,2.8,4.75,[
    {"t":"핵심 결과","size":15,"color":BLUE_L,"bold":True,"space_after":10},
    {"t":"고장 탐지 recall","size":12,"color":SLATE,"space_after":2},
    {"t":"78.8% → 90.4%","size":19,"color":GREEN2,"bold":True,"space_after":12},
    {"t":"정확도","size":12,"color":SLATE,"space_after":2},
    {"t":"92.1% → 95.9%","size":18,"color":GREEN2,"bold":True,"space_after":14},
    {"t":"PSI가 0.2를 넘은 t*=6,700에서 자동 재학습 트리거 → 초록(운영)이 빨강(미재학습) 위로 회복","size":11.5,"color":SLATE,"line_spacing":1.3},
])
txt(s,0.5,6.55,8.75,0.4,"좌: 입력 드리프트 지표 PSI(임계 0.2 돌파) · 우: 고장 탐지 recall — 빨강=미재학습 vs 초록=운영(재학습)",size=10.5,color=GRAY)
pagenum(s,13)

# ============================================================ Slide 14 데모 화면 (SPC & RF)
s=blank(); title_bar(s,"데모 — SPC 실시간 감시 & RF 고장 진단","LIVE DEMO",accent=PURPLE)
add_img(s,"shots/02_spc_drift.png",0.45,1.65,6.0)
add_img(s,"shots/03_rf_diagnosis.png",6.85,1.65,6.0)
txt(s,0.45,5.35,6.0,0.9,[
    {"t":"① Layer 1 · SPC 관리도","size":13,"color":AMBER,"bold":True,"space_after":2},
    {"t":"5개 센서 실시간 감시 — 드리프트 구간에서 토크가 관리한계를 이탈(빨간 점). 라벨 없이 이상 감지 + 입력 드리프트 감시.","size":11.5,"color":DARK,"line_spacing":1.15},
])
txt(s,6.85,5.35,6.0,0.9,[
    {"t":"② Layer 2 · Random Forest 진단","size":13,"color":PURPLE,"bold":True,"space_after":2},
    {"t":"고장 모드별 진단 분포(실제 vs 모델) + 최근 진단 내역(확신도·일치 여부). \"무슨 고장인지\"까지 자동 판정.","size":11.5,"color":DARK,"line_spacing":1.15},
])
pagenum(s,14)

# ============================================================ Slide 15 결론 & 향후
s=blank(); title_bar(s,"결론 & 향후 계획","CONCLUSION",accent=BLUE)
add_img(s,"shots/05_business.png",0.5,1.55,5.75)
txt(s,0.5,5.5,5.75,0.4,"④ 비즈니스 가치 탭 — 재학습 전/후 효과 실측",size=11,color=GRAY)
rect(s,6.6,1.55,6.15,2.25,NAVY,rounded=True)
txt(s,6.85,1.7,5.7,2.0,[
    {"t":"결론","size":15,"color":BLUE_L,"bold":True,"space_after":6},
    {"t":"이중 레이어(SPC+RF)가 감지·진단·자기유지를 동시에 수행하고, 두 레이어가 서로를 검증해 모델이 조용히 늙는 것을 막는다. SPC는 1차 감시이자 ML 모델의 드리프트 감시자 역할을 겸한다.","size":12.5,"color":WHITE,"line_spacing":1.25},
])
rect(s,6.6,4.0,6.15,2.5,GREEN_L,rounded=True)
txt(s,6.85,4.15,5.7,2.25,[
    {"t":"향후 계획","size":15,"color":GREEN,"bold":True,"space_after":6},
    {"t":"• 라벨 지연을 반영한 backfill·지연 평가 추가","size":12.5,"color":DARK,"space_after":4},
    {"t":"• 추가 센서·고장 모드로 확장","size":12.5,"color":DARK,"space_after":4},
    {"t":"• Champion/Challenger 자동 승급 + Shadow 배포","size":12.5,"color":DARK,"space_after":4},
    {"t":"• 실시간 스트리밍(Kafka) 연동","size":12.5,"color":DARK},
])
pagenum(s,15)

_out = HERE / "발표_설비고장_SPC이중레이어_MLOps.pptx"
prs.save(str(_out))
print("저장 완료:", _out)
print("슬라이드 수:", len(prs.slides._sldIdLst))
